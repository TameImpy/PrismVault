"""One-off template edit: add the provenance markers to the Appendix slide (#156).

The deck template is the single source of truth for styling (PRD #131), so a
new marker has to be *authored into the template* rather than drawn by
deck_builder at build time. This script does that authoring once and is kept in
the repo so the edit is reproducible and reviewable — the .pptx it produces is
a binary, and a diff of it says nothing.

What it does: clones the Methodology slide's body text box (so the new box
inherits Barlow, the colour and the bullet styling verbatim rather than any of
it being invented here), drops it on the empty Appendix slide, and fills it
with a static heading plus one `[PROVENANCE_n]` marker paragraph per data
section. Only the font *size* is changed, down to footnote scale, because a
source line is reference matter and has to fit beside five others.

Idempotent: re-running replaces the box rather than adding a second one.

    python3 scripts/add_provenance_markers.py
"""
import copy
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Pt

PROJECT_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from src.deck_builder import TEMPLATE_PATH, PROVENANCE_TILES  # noqa: E402

APPENDIX_TITLE = "Appendix"
METHODOLOGY_TITLE = "Methodology"
METHODOLOGY_BODY_START = "Audience segments"

HEADING = "Every figure in this deck traces back to one of these sources."

# Footnote scale. The box has to hold six source lines without crowding the
# slide; everything else about the look comes from the cloned box.
FOOTNOTE_SIZE = Pt(10)

# Slide margins match the rest of the deck; the box is widened to the facing
# margin so a long source line wraps twice rather than four times.
BOX_LEFT = Emu(631075)
BOX_TOP = Emu(984251)
BOX_WIDTH = Emu(7881850)
BOX_HEIGHT = Emu(2800000)


def _slide_with_title(prs, title):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() == title:
                return slide
    raise SystemExit("No slide titled %r in the template." % title)


def _methodology_body(prs):
    slide = _slide_with_title(prs, METHODOLOGY_TITLE)
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith(METHODOLOGY_BODY_START):
            return shape
    raise SystemExit("Could not find the Methodology body text box to clone.")


def _clear_extra_paragraphs(text_frame, keep):
    """Trim the cloned box down to `keep` paragraphs, template styling intact."""
    paragraphs = text_frame.paragraphs
    for para in paragraphs[keep:]:
        para._p.getparent().remove(para._p)


def _set_single_run(para, text, size):
    """Collapse a paragraph to one run carrying `text` at `size`."""
    runs = para.runs
    for run in runs[1:]:
        run._r.getparent().remove(run._r)
    runs[0].text = text
    runs[0].font.size = size


def build(template_path=TEMPLATE_PATH):
    prs = Presentation(template_path)
    appendix = _slide_with_title(prs, APPENDIX_TITLE)

    # Idempotence: drop any provenance box a previous run left behind.
    for shape in list(appendix.shapes):
        if shape.has_text_frame and "[PROVENANCE_1]" in shape.text_frame.text:
            shape._element.getparent().remove(shape._element)

    source = _methodology_body(prs)
    clone = copy.deepcopy(source._element)
    appendix.shapes._spTree.append(clone)

    box = [s for s in appendix.shapes if s._element is clone][0]
    box.left, box.top, box.width, box.height = BOX_LEFT, BOX_TOP, BOX_WIDTH, BOX_HEIGHT

    lines = [HEADING] + ["[PROVENANCE_%d]" % (i + 1) for i in range(PROVENANCE_TILES)]
    text_frame = box.text_frame

    # The clone brings more paragraphs than we need and fewer than we might —
    # size it to the line count by trimming, then cloning its last paragraph.
    _clear_extra_paragraphs(text_frame, 1)
    template_para = text_frame.paragraphs[0]._p
    for _ in range(len(lines) - 1):
        text_frame.paragraphs[0]._p.getparent().append(copy.deepcopy(template_para))

    for para, text in zip(text_frame.paragraphs, lines):
        _set_single_run(para, text, FOOTNOTE_SIZE)

    prs.save(template_path)
    print("Added %d provenance markers to the Appendix slide of %s"
          % (PROVENANCE_TILES, template_path))


if __name__ == "__main__":
    build()
