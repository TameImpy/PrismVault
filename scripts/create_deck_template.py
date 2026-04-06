"""
One-time script to create the Prism Plan content template from the master deck.

Extracts background images and logos from the Prism Story Deck Master,
then creates a new .pptx with 5 reusable slide layouts containing proper
named placeholders for dynamic content.

Usage:
    python3 scripts/create_deck_template.py
"""
import sys
import os
import io

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from lxml import etree

MASTER_PATH = os.path.expanduser("~/Desktop/Prism Story Deck Master.pptx")
OUTPUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "data", "templates", "prism_plan_template.pptx")

# Slide dimensions matching master
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500

# Brand colours
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xC0, 0xC0, 0xC0)
CYAN_ACCENT = RGBColor(0x1F, 0x89, 0xDF)

# Font sizes in EMU (Pt * 12700)
TITLE_SIZE = Pt(40)
SUBTITLE_SIZE = Pt(16)
HEADING_SIZE = Pt(24)
BODY_SIZE = Pt(14)
STAT_VALUE_SIZE = Pt(40)
STAT_LABEL_SIZE = Pt(11)
SMALL_SIZE = Pt(10)
CTA_SIZE = Pt(14)


def _extract_image_blob(shape):
    """Extract the image blob from a picture shape."""
    return shape.image.blob, shape.image.content_type


def _add_background_image(slide, blob, content_type):
    """Add a full-bleed background image to a slide."""
    stream = io.BytesIO(blob)
    pic = slide.shapes.add_picture(stream, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    # Move to back
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _add_logo(slide, blob, content_type, left, top, width, height):
    """Add a logo image to a slide."""
    stream = io.BytesIO(blob)
    slide.shapes.add_picture(stream, left, top, width, height)


def _add_textbox(slide, left, top, width, height, text="", font_name="Barlow ExtraBold",
                 font_size=HEADING_SIZE, font_color=WHITE, alignment=PP_ALIGN.LEFT,
                 bold=None, word_wrap=True):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    # Enable auto-fit (shrink text)
    txBody = tf._txBody
    bodyPr = txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
    if bodyPr is not None:
        normAutofit = etree.SubElement(
            bodyPr,
            '{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit'
        )
        normAutofit.set('fontScale', '62500')  # Min ~62.5% of original size
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    if bold is not None:
        run.font.bold = bold
    return txBox


def create_template():
    """Create the Prism Plan content template."""
    # Load master to extract assets
    master = Presentation(MASTER_PATH)

    # Extract background from slide 1 (dark branded bg)
    slide1 = master.slides[0]
    bg_blob = None
    bg_ct = None
    logo1_data = None
    logo2_data = None

    for shape in slide1.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if shape.width > 8000000 and shape.height > 4000000:
                bg_blob, bg_ct = _extract_image_blob(shape)
            elif shape.left > 7000000 and shape.top > 4000000:
                logo2_data = (_extract_image_blob(shape), shape.left, shape.top, shape.width, shape.height)
            elif shape.left > 5000000 and shape.top > 4000000:
                logo1_data = (_extract_image_blob(shape), shape.left, shape.top, shape.width, shape.height)

    # Also extract a different background for variety (slide 4 — stats-style)
    slide4 = master.slides[3]
    bg2_blob = None
    bg2_ct = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.width > 8000000:
            bg2_blob, bg2_ct = _extract_image_blob(shape)
            break

    # Extract background from slide 9 (transition style)
    slide9 = master.slides[8]
    bg3_blob = None
    bg3_ct = None
    for shape in slide9.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.width > 8000000:
            bg3_blob, bg3_ct = _extract_image_blob(shape)
            break

    # Create new presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use the blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # ===== SLIDE 1: TITLE =====
    slide = prs.slides.add_slide(blank_layout)
    if bg_blob:
        _add_background_image(slide, bg_blob, bg_ct)

    # Topic (large)
    _add_textbox(slide, Emu(631075), Emu(669941), Emu(7408200), Emu(1468800),
                 text="[TOPIC]", font_name="Barlow ExtraBold", font_size=Pt(50),
                 font_color=WHITE)
    # Advertiser
    _add_textbox(slide, Emu(752112), Emu(2649444), Emu(7408200), Emu(431100),
                 text="[ADVERTISER]", font_name="Barlow Medium", font_size=Pt(16),
                 font_color=LIGHT_GREY)
    # KPI
    _add_textbox(slide, Emu(752112), Emu(3200000), Emu(7408200), Emu(300000),
                 text="KPI: [KPI]", font_name="Barlow", font_size=Pt(14),
                 font_color=CYAN_ACCENT)

    if logo1_data:
        (blob, ct), l, t, w, h = logo1_data
        _add_logo(slide, blob, ct, l, t, w, h)
    if logo2_data:
        (blob, ct), l, t, w, h = logo2_data
        _add_logo(slide, blob, ct, l, t, w, h)

    # ===== SLIDE 2: STATS =====
    slide = prs.slides.add_slide(blank_layout)
    if bg2_blob:
        _add_background_image(slide, bg2_blob, bg2_ct)
    elif bg_blob:
        _add_background_image(slide, bg_blob, bg_ct)

    # Section heading
    _add_textbox(slide, Emu(631075), Emu(280000), Emu(7884000), Emu(500000),
                 text="At a Glance", font_name="Barlow ExtraBold", font_size=Pt(32),
                 font_color=WHITE)

    # 6 stat boxes in a 3x2 grid
    col_positions = [631075, 3300000, 5968925]
    row_positions = [950000, 2900000]
    for row_idx, row_top in enumerate(row_positions):
        for col_idx, col_left in enumerate(col_positions):
            # Value (big number)
            _add_textbox(slide, Emu(col_left), Emu(row_top), Emu(2400000), Emu(700000),
                         text="[VALUE]", font_name="Barlow ExtraBold", font_size=STAT_VALUE_SIZE,
                         font_color=WHITE)
            # Label
            _add_textbox(slide, Emu(col_left), Emu(row_top + 700000), Emu(2400000), Emu(450000),
                         text="[LABEL]", font_name="Barlow", font_size=STAT_LABEL_SIZE,
                         font_color=LIGHT_GREY)

    if logo2_data:
        (blob, ct), l, t, w, h = logo2_data
        _add_logo(slide, blob, ct, l, t, w, h)

    # ===== SLIDE 3: TWO-COLUMN BODY =====
    slide = prs.slides.add_slide(blank_layout)
    if bg_blob:
        _add_background_image(slide, bg_blob, bg_ct)

    # Left heading
    _add_textbox(slide, Emu(631075), Emu(350000), Emu(3800000), Emu(500000),
                 text="Advertiser Overview", font_name="Barlow ExtraBold", font_size=Pt(22),
                 font_color=WHITE)
    # Left body
    _add_textbox(slide, Emu(631075), Emu(900000), Emu(3800000), Emu(3600000),
                 text="[LEFT_BODY]", font_name="Barlow", font_size=BODY_SIZE,
                 font_color=LIGHT_GREY)
    # Right heading
    _add_textbox(slide, Emu(4800000), Emu(350000), Emu(3800000), Emu(500000),
                 text="Editorial Insights", font_name="Barlow ExtraBold", font_size=Pt(22),
                 font_color=WHITE)
    # Right body
    _add_textbox(slide, Emu(4800000), Emu(900000), Emu(3800000), Emu(3600000),
                 text="[RIGHT_BODY]", font_name="Barlow", font_size=BODY_SIZE,
                 font_color=LIGHT_GREY)

    if logo2_data:
        (blob, ct), l, t, w, h = logo2_data
        _add_logo(slide, blob, ct, l, t, w, h)

    # ===== SLIDE 4: BODY (Messaging & Tone) =====
    slide = prs.slides.add_slide(blank_layout)
    if bg3_blob:
        _add_background_image(slide, bg3_blob, bg3_ct)
    elif bg_blob:
        _add_background_image(slide, bg_blob, bg_ct)

    # Heading
    _add_textbox(slide, Emu(631075), Emu(350000), Emu(7884000), Emu(600000),
                 text="Messaging & Tone", font_name="Barlow ExtraBold", font_size=Pt(32),
                 font_color=WHITE)
    # Body area
    _add_textbox(slide, Emu(631075), Emu(1050000), Emu(7884000), Emu(3500000),
                 text="[BODY]", font_name="Barlow", font_size=BODY_SIZE,
                 font_color=LIGHT_GREY)

    if logo2_data:
        (blob, ct), l, t, w, h = logo2_data
        _add_logo(slide, blob, ct, l, t, w, h)

    # ===== SLIDE 5: CLOSING/CTA =====
    slide = prs.slides.add_slide(blank_layout)
    if bg2_blob:
        _add_background_image(slide, bg2_blob, bg2_ct)
    elif bg_blob:
        _add_background_image(slide, bg_blob, bg_ct)

    # Heading
    _add_textbox(slide, Emu(631075), Emu(350000), Emu(7884000), Emu(600000),
                 text="Recommended Products", font_name="Barlow ExtraBold", font_size=Pt(32),
                 font_color=WHITE)
    # Products area
    _add_textbox(slide, Emu(631075), Emu(1050000), Emu(7884000), Emu(2400000),
                 text="[PRODUCTS]", font_name="Barlow", font_size=BODY_SIZE,
                 font_color=LIGHT_GREY)
    # CTA
    _add_textbox(slide, Emu(631075), Emu(3700000), Emu(7884000), Emu(800000),
                 text="[CTA]", font_name="Barlow Medium", font_size=CTA_SIZE,
                 font_color=CYAN_ACCENT, alignment=PP_ALIGN.CENTER)

    if logo1_data:
        (blob, ct), l, t, w, h = logo1_data
        _add_logo(slide, blob, ct, l, t, w, h)
    if logo2_data:
        (blob, ct), l, t, w, h = logo2_data
        _add_logo(slide, blob, ct, l, t, w, h)

    # Save
    prs.save(OUTPUT_PATH)
    print("Template saved to: %s" % OUTPUT_PATH)
    print("Slides: %d" % len(prs.slides))


if __name__ == "__main__":
    create_template()
