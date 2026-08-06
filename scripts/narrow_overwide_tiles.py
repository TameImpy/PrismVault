"""One-off template edit: narrow the tile text boxes to their columns (#162).

Reported as "some exported slides have text overlapping". The cause is
geometry, not content. Several marker boxes in the official template are drawn
wider than the column they belong to — `[INSIGHT_2_STAT]` is a 5.34" box at
x=4.88" on a 10" slide, so it extends 0.22" past the right edge — and every box
is `spAutoFit`. PowerPoint wraps text at the *box* width and grows the shape
downward, so a value long enough to use the box's full width is drawn across
the neighbouring tile, or over the Prism logo. No amount of shortening the text
fixes that: on slide 5 the only string guaranteed to stay inside tile 1 with
the box as drawn is about four words long, which would gut the insight tiles.

So the boxes are narrowed to the column each tile actually occupies. Nothing
else changes: no font, size, colour, position or anchoring is touched, and the
text is left-aligned and top-anchored, so for every value short enough to fit
today the rendering is byte-for-byte what it was. The edit only bites where
text used to spill.

The widths below are the columns, measured from the template itself — the gap
to the next tile to the right, or to the slide edge, or to the logo artwork.
`src.tile_fit.TILE_BUDGETS` holds the matching usable widths and
`tests/test_tile_fit.py` asserts the two agree, so the pair cannot drift.

Idempotent: the widths are absolute, so re-running after a template revision
re-applies them rather than shrinking anything twice.

    python3 scripts/narrow_overwide_tiles.py
"""
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches

PROJECT_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from src.deck_builder import TEMPLATE_PATH  # noqa: E402

MARKER_RE = re.compile(r"\[[A-Z0-9_]+\]")

# Marker family -> box width in inches.
#
# Products and segments: 4.80" is the gap from the left column (x=0.35"/0.22")
# to the right column (x=5.15"/5.18"), and the right column has the same width
# again before the slide edge. The third product tile sits between them and
# takes the same width so all three wrap alike.
#
# Insights: 3.74" for the supporting line. The column is 4.53" wide, but the
# third tile sits under the Prism logo at x=6.41", and 2.57" + 3.74" clears it.
# All three tiles take that width rather than only the one that needs it —
# three tiles drawn identically have to wrap identically. The 40pt figure above
# it is a single short string and keeps the full 4.53" column.
#
# The 2.62" name and format boxes are already inside their columns and are the
# deliberate width of the label beneath a figure; they are left alone.
COLUMN_WIDTHS = {
    "[PRODUCT_N_CTR]": 4.80,
    "[PRODUCT_N_VIEW]": 4.80,
    "[SEGMENT_N_REACH]": 4.80,
    "[INSIGHT_N]": 4.53,
    "[INSIGHT_N_STAT]": 3.74,
}


def family(marker):
    return re.sub(r"_\d+", "_N", marker)


def narrow(path=TEMPLATE_PATH):
    prs = Presentation(path)
    changed = []

    for index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            markers = MARKER_RE.findall(shape.text_frame.text)
            if not markers:
                continue
            target = COLUMN_WIDTHS.get(family(markers[0]))
            if target is None:
                continue
            was = shape.width
            shape.width = Inches(target)
            if shape.width != was:
                changed.append("slide %d  %-20s %-12s %.2f\" -> %.2f\""
                               % (index, markers[0], shape.name,
                                  was / 914400.0, target))

    prs.save(path)
    return changed


if __name__ == "__main__":
    for line in narrow():
        print(line)
    print("template written: %s" % os.path.relpath(TEMPLATE_PATH, PROJECT_ROOT))
