"""
Deck builder. Populates the official 8-slide template with the structured brief
payload and returns the finished .pptx as a BytesIO buffer.

The template is the single source of truth for styling (PRD #131): this module
injects text into the template's own marker runs and never touches a font, size
or colour. Change the look by editing the template, not this file.
"""
import io
import os

from pptx import Presentation
from pptx.oxml.ns import qn

from src.font_embed import embed_fonts
from src.provenance import format_provenance_row
from src.slide_content import (
    INSIGHT_TILES,
    PRODUCT_TILES,
    PROVENANCE_TILES,
    SEGMENT_TILES,
    build_product_rows,
    build_segment_rows,
)
from src.tile_fit import MARKER_RE, fit_to_tile

TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "data", "templates", "prism_plan_template_export.pptx",
)

# The Historical Insights slide is dropped whole when no research matched.
INSIGHT_SLIDE_MARKER = "[INSIGHT_1]"

_XML_SPACE = "{http://www.w3.org/XML/1998/namespace}space"


def build_deck(slide_content, advertiser):
    """Build the deck from the brief's structured payload.

    Args:
        slide_content: the deck payload — ``advertiser_overview`` and
            ``insights`` from the LLM step, plus the brief's own
            ``audience_segments`` and ``format_recommendations``.
        advertiser: the advertiser name, as the brief was run for.

    Returns:
        io.BytesIO containing the finished .pptx, with Barlow embedded.
    """
    prs = Presentation(TEMPLATE_PATH)

    if not slide_content.get("insights"):
        _drop_slide_with_marker(prs, INSIGHT_SLIDE_MARKER)

    _fill_markers(prs, _build_fields(slide_content, advertiser))

    buf = io.BytesIO()
    prs.save(buf)
    # Embed Barlow so the deck renders on-brand on a machine without it installed.
    return embed_fonts(buf)


def _build_fields(slide_content, advertiser):
    """Map the payload onto every marker the template defines.

    Every marker gets an entry — tiles with no data behind them resolve to ""
    and are blanked, so no `[MARKER]` can survive into a client-facing deck.

    Values leave here fitted to their tile (`src/tile_fit.py`): one choke point,
    so a marker added later cannot escape the guard that stops long content
    landing on the tile below it (#162).
    """
    fields = {
        "[ADVERTISER_NAME]": advertiser,
        "[ADVERTISER_OVERVIEW]": slide_content.get("advertiser_overview"),
    }

    products = build_product_rows(slide_content.get("format_recommendations"))
    for i in range(PRODUCT_TILES):
        row = products[i] if i < len(products) else {}
        fields["[PRODUCT_%d_FORMAT]" % (i + 1)] = row.get("format")
        fields["[PRODUCT_%d_CTR]" % (i + 1)] = row.get("ctr")
        fields["[PRODUCT_%d_VIEW]" % (i + 1)] = row.get("viewability")

    segments = build_segment_rows(slide_content.get("audience_segments"))
    for i in range(SEGMENT_TILES):
        row = segments[i] if i < len(segments) else {}
        fields["[SEGMENT_%d_NAME]" % (i + 1)] = row.get("name")
        fields["[SEGMENT_%d_REACH]" % (i + 1)] = row.get("reach")

    insights = slide_content.get("insights") or []
    for i in range(INSIGHT_TILES):
        row = insights[i] if i < len(insights) else {}
        # The hero line of each tile is the figure (40pt ExtraBold, mirroring
        # the product CTR tile); the supporting phrase sits beneath it.
        fields["[INSIGHT_%d]" % (i + 1)] = row.get("stat")
        fields["[INSIGHT_%d_STAT]" % (i + 1)] = row.get("text")

    # The appendix names the source, date, coverage and period behind each data
    # section the brief drew on (#156), one line per section. The lines are the
    # brief's own — placed verbatim, never regenerated for the deck.
    provenance = slide_content.get("provenance") or []
    for i in range(PROVENANCE_TILES):
        entry = provenance[i] if i < len(provenance) else None
        fields["[PROVENANCE_%d]" % (i + 1)] = format_provenance_row(entry)

    return {marker: fit_to_tile(marker, value) for marker, value in fields.items()}


def _drop_slide_with_marker(prs, marker):
    """Remove the slide carrying `marker`, along with its relationship."""
    for index, slide in enumerate(prs.slides):
        if marker in _slide_text(slide):
            sldId = list(prs.slides._sldIdLst)[index]
            prs.part.drop_rel(sldId.rId)
            prs.slides._sldIdLst.remove(sldId)
            return


def _slide_text(slide):
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def _fill_markers(prs, fields):
    """Replace every marker in the deck with its value, in place.

    Text is written into the template's own runs, so each populated run keeps
    the placeholder's font, size and colour untouched.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _fill_paragraph(para, fields)


def _fill_paragraph(para, fields):
    markers = MARKER_RE.findall(para.text)
    if not markers:
        return

    # A marker occupies a single run in the template. If a template revision
    # ever splits one across runs, collapse the paragraph into its first run so
    # the fill still lands — that run's formatting is the placeholder's.
    if any(not any(m in r.text for r in para.runs) for m in markers):
        _collapse_runs(para)

    for run in para.runs:
        text = run.text
        for marker in markers:
            if marker in text:
                text = text.replace(marker, fields.get(marker, ""))
        if text != run.text:
            _set_run_text(run, text)

    _drop_empty_runs(para)


def _collapse_runs(para):
    """Merge a paragraph's runs into the first one, keeping its formatting."""
    runs = para.runs
    if len(runs) < 2:
        return
    _set_run_text(runs[0], para.text)
    for run in runs[1:]:
        run._r.getparent().remove(run._r)


def _set_run_text(run, text):
    run.text = text
    # Leading/trailing spaces are meaningful here ("[ADVERTISER_NAME] Brief").
    t = run._r.find(qn("a:t"))
    if t is not None:
        t.set(_XML_SPACE, "preserve")


def _drop_empty_runs(para):
    """Remove runs left with no text — an empty <a:t/> reads as corrupt."""
    for run in para.runs:
        if not run.text:
            run._r.getparent().remove(run._r)
