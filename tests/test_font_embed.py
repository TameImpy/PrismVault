"""
Tests for the font embedder at src/font_embed.py.

No mocks — embeds into a real .pptx and inspects the resulting OPC package.
"""
import io
import os
import sys
import zipfile

PROJECT_ROOT = os.path.join(os.path.dirname(__file__), "..")
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

import pytest
from lxml import etree
from pptx import Presentation

from src.font_embed import (
    BARLOW_FONTS,
    FONTS_DIR,
    FONT_CONTENT_TYPE,
    FONT_REL_TYPE,
    embed_fonts,
)

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PR_NS = "http://schemas.openxmlformats.org/package/2006/relationships"

TEMPLATE_PATH = os.path.join(PROJECT_ROOT, "data", "templates", "prism_plan_template_export.pptx")


def _template_bytes():
    with open(TEMPLATE_PATH, "rb") as f:
        return io.BytesIO(f.read())


@pytest.fixture(scope="module")
def embedded():
    """The official template with the Barlow weights embedded."""
    return embed_fonts(_template_bytes())


def _zip(buf):
    buf.seek(0)
    return zipfile.ZipFile(io.BytesIO(buf.read()))


def _presentation_xml(zf):
    return etree.fromstring(zf.read("ppt/presentation.xml"))


def _presentation_rels(zf):
    """{rId: (type, target)} for ppt/presentation.xml."""
    root = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    return {
        rel.get("Id"): (rel.get("Type"), rel.get("Target"))
        for rel in root.findall("{%s}Relationship" % PR_NS)
    }


def _embedded_fonts(pres_xml):
    """[(typeface, {slot: rId})] from the embedded-font list."""
    lst = pres_xml.find("{%s}embeddedFontLst" % P_NS)
    if lst is None:
        return []
    out = []
    for entry in lst.findall("{%s}embeddedFont" % P_NS):
        font = entry.find("{%s}font" % P_NS)
        slots = {}
        for child in entry:
            tag = etree.QName(child).localname
            if tag == "font":
                continue
            slots[tag] = child.get("{%s}id" % R_NS)
        out.append((font.get("typeface"), slots))
    return out


# --- The embedder in isolation -------------------------------------------------


def test_embeds_the_three_barlow_weights(embedded):
    """Regular, ExtraBold and ExtraLight are listed under their template typefaces."""
    fonts = _embedded_fonts(_presentation_xml(_zip(embedded)))

    assert [typeface for typeface, _ in fonts] == [
        "Barlow",
        "Barlow ExtraBold",
        "Barlow ExtraLight",
    ]


def test_each_weight_occupies_the_regular_slot(embedded):
    """Each Barlow weight is its own font family, so each fills its family's regular slot."""
    fonts = _embedded_fonts(_presentation_xml(_zip(embedded)))

    for typeface, slots in fonts:
        assert list(slots) == ["regular"], typeface


def test_presentation_is_flagged_to_embed_fonts(embedded):
    pres = _presentation_xml(_zip(embedded))

    assert pres.get("embedTrueTypeFonts") == "1"


def test_full_fonts_are_embedded_not_subsets(embedded):
    """saveSubsetFonts is cleared — we embed complete fonts so recipients can edit."""
    pres = _presentation_xml(_zip(embedded))

    assert pres.get("saveSubsetFonts") == "0"


def test_font_parts_hold_the_committed_ttf_bytes(embedded):
    """Each font part is the raw .ttf from data/fonts, byte for byte."""
    zf = _zip(embedded)
    rels = _presentation_rels(zf)
    fonts = _embedded_fonts(_presentation_xml(zf))

    for (typeface, slots), spec in zip(fonts, BARLOW_FONTS):
        rel_type, target = rels[slots["regular"]]
        assert rel_type == FONT_REL_TYPE
        with open(os.path.join(FONTS_DIR, spec.filename), "rb") as f:
            assert zf.read("ppt/" + target) == f.read(), typeface


def test_font_relationships_resolve_to_parts_that_exist(embedded):
    zf = _zip(embedded)
    rels = _presentation_rels(zf)
    names = set(zf.namelist())

    font_rels = [t for (rt, t) in rels.values() if rt == FONT_REL_TYPE]
    assert len(font_rels) == len(BARLOW_FONTS)
    for target in font_rels:
        assert "ppt/" + target in names


def test_font_parts_carry_the_font_content_type(embedded):
    """A Default mapping for the fntdata extension — what PowerPoint itself writes."""
    zf = _zip(embedded)
    types = etree.fromstring(zf.read("[Content_Types].xml"))

    defaults = {
        d.get("Extension").lower(): d.get("ContentType")
        for d in types.findall("{%s}Default" % CT_NS)
    }
    assert defaults.get("fntdata") == FONT_CONTENT_TYPE


def test_font_relationship_ids_do_not_collide(embedded):
    """New rIds are allocated above the template's existing ones."""
    zf = _zip(embedded)
    root = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
    ids = [rel.get("Id") for rel in root.findall("{%s}Relationship" % PR_NS)]

    assert len(ids) == len(set(ids))


def test_embedded_font_list_is_in_schema_order(embedded):
    """CT_Presentation requires embeddedFontLst after notesSz and before defaultTextStyle."""
    pres = _presentation_xml(_zip(embedded))
    order = [etree.QName(child).localname for child in pres]

    idx = order.index("embeddedFontLst")
    for earlier in ("sldMasterIdLst", "sldIdLst", "sldSz", "notesSz"):
        if earlier in order:
            assert order.index(earlier) < idx, earlier
    for later in ("custShowLst", "defaultTextStyle", "extLst"):
        if later in order:
            assert idx < order.index(later), later


def test_embedding_is_idempotent(embedded):
    """Re-embedding replaces the font parts rather than duplicating them."""
    twice = embed_fonts(embedded)
    zf = _zip(twice)

    fonts = _embedded_fonts(_presentation_xml(zf))
    assert len(fonts) == len(BARLOW_FONTS)
    assert len([n for n in zf.namelist() if n.endswith(".fntdata")]) == len(BARLOW_FONTS)


def test_original_parts_are_preserved(embedded):
    """Embedding only adds parts — nothing from the source package is dropped."""
    before = set(_zip(_template_bytes()).namelist())
    after = set(_zip(embedded).namelist())

    assert before <= after
    assert {n for n in after - before} == {
        "ppt/fonts/font%d.fntdata" % (i + 1) for i in range(len(BARLOW_FONTS))
    }


def test_content_types_first_in_the_archive(embedded):
    """OPC requires the content-types item to be the first entry in the zip."""
    assert _zip(embedded).namelist()[0] == "[Content_Types].xml"


def test_output_still_opens_as_a_presentation(embedded):
    """The slides survive the repackaging untouched."""
    source = Presentation(_template_bytes())
    result = Presentation(embedded)

    assert len(result.slides) == len(source.slides)
    assert [
        " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        for slide in result.slides
    ] == [
        " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        for slide in source.slides
    ]


def test_accepts_raw_bytes(embedded):
    with open(TEMPLATE_PATH, "rb") as f:
        result = embed_fonts(f.read())

    assert _presentation_xml(_zip(result)).get("embedTrueTypeFonts") == "1"


def test_missing_font_file_is_reported(monkeypatch):
    monkeypatch.setattr("src.font_embed.FONTS_DIR", os.path.join(PROJECT_ROOT, "data", "nope"))

    with pytest.raises(FileNotFoundError):
        embed_fonts(_template_bytes())


# --- End to end, through the deck build ---------------------------------------


def test_built_deck_embeds_the_brand_fonts():
    from tests.test_deck_builder import SAMPLE_CONTENT
    from src.deck_builder import build_deck

    zf = _zip(build_deck(SAMPLE_CONTENT, "gut health", "Yakult"))
    pres = _presentation_xml(zf)

    assert pres.get("embedTrueTypeFonts") == "1"
    assert [typeface for typeface, _ in _embedded_fonts(pres)] == [
        "Barlow",
        "Barlow ExtraBold",
        "Barlow ExtraLight",
    ]
    assert sorted(n for n in zf.namelist() if n.endswith(".fntdata")) == [
        "ppt/fonts/font1.fntdata",
        "ppt/fonts/font2.fntdata",
        "ppt/fonts/font3.fntdata",
    ]


def test_built_deck_is_structurally_intact():
    """Every relationship in the built deck resolves to a part that exists, and every
    part is typed — the two package-level rules PowerPoint fails a file's repair check on."""
    from tests.test_deck_builder import SAMPLE_CONTENT
    from src.deck_builder import build_deck

    zf = _zip(build_deck(SAMPLE_CONTENT, "gut health", "Yakult"))
    names = set(zf.namelist())

    defaults = {}
    overrides = {}
    types = etree.fromstring(zf.read("[Content_Types].xml"))
    for d in types.findall("{%s}Default" % CT_NS):
        defaults[d.get("Extension").lower()] = d.get("ContentType")
    for o in types.findall("{%s}Override" % CT_NS):
        overrides[o.get("PartName")] = o.get("ContentType")

    untyped = [
        n for n in names
        if "/" + n not in overrides and n.rsplit(".", 1)[-1].lower() not in defaults
    ]
    assert untyped == []

    dangling = []
    for name in names:
        if not name.endswith(".rels"):
            continue
        base = name.replace("_rels/", "").rsplit("/", 1)[0] if "/" in name else ""
        base = name[: name.index("_rels/")]
        for rel in etree.fromstring(zf.read(name)).findall("{%s}Relationship" % PR_NS):
            target = rel.get("Target")
            if rel.get("TargetMode") == "External" or target.startswith(("http", "/")):
                continue
            resolved = os.path.normpath(os.path.join(base, target)).replace(os.sep, "/")
            if resolved not in names:
                dangling.append((name, target))
    assert dangling == []
