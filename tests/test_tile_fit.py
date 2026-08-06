"""
Tests for the tile fitter at src/tile_fit.py (#162).

Two halves, and the second is the one that matters.

The first checks the fitting itself — measuring a string in Barlow, wrapping it
the way PowerPoint would, and shortening it when it will not fit.

The second pins every budget in ``TILE_BUDGETS`` against the template it was
read from: the typeface and point size are the template's own, no budget
overlaps another or runs off the slide, and every marker the template defines
either has a budget or is a documented exemption. That is what stops the table
becoming a set of numbers nobody can re-derive after a template revision.
"""
import sys
import os

PROJECT_ROOT = os.path.join(os.path.dirname(__file__), "..")
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

import csv
import re

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from src.deck_builder import TEMPLATE_PATH
from src.tile_fit import (
    ELLIPSIS,
    EXEMPT_MARKERS,
    MARKER_RE,
    TILE_BUDGETS,
    budget_for,
    fit_to_tile,
    line_height,
    inches,
    marker_family,
    rendered_lines,
    text_width,
    wrapping_width,
)
from tests.tile_geometry import artwork_rects, overlaps, text_origin


# ---------------------------------------------------------------------------
# Measuring and wrapping
# ---------------------------------------------------------------------------


def test_text_width_grows_with_the_string():
    assert text_width("mm", "Barlow", 26) > text_width("m", "Barlow", 26)


def test_text_width_scales_with_the_point_size():
    assert text_width("Homepride", "Barlow", 40) == pytest.approx(
        2 * text_width("Homepride", "Barlow", 20), rel=1e-6
    )


def test_text_width_uses_real_barlow_advances_not_a_flat_estimate():
    """A narrow glyph must measure narrower than a wide one — the half-the-point-
    size estimate used elsewhere cannot tell these apart."""
    assert text_width("i" * 10, "Barlow", 26) < text_width("W" * 10, "Barlow", 26)


def test_unknown_typeface_still_measures():
    """A template revision naming a weight we hold no file for must not crash the
    download — it falls back to Barlow regular."""
    assert text_width("Homepride", "Barlow Semi Whatever", 26) > 0


def test_wrapping_breaks_on_words():
    """Wrapping splits between words and loses none of them."""
    lines = rendered_lines("one two three four five six", 0.8, "Barlow", 11)
    assert len(lines) > 1
    assert " ".join(lines) == "one two three four five six"


def test_wrapping_keeps_a_short_string_on_one_line():
    assert rendered_lines("68%", 4.0, "Barlow", 26) == ["68%"]


def test_line_height_scales_with_the_point_size():
    assert line_height("Barlow", 26) == pytest.approx(2 * line_height("Barlow", 13))


# ---------------------------------------------------------------------------
# Fitting
# ---------------------------------------------------------------------------


def test_text_that_fits_is_returned_untouched():
    assert fit_to_tile("[SEGMENT_1_NAME]", "Confident bakers") == "Confident bakers"


def test_text_too_long_is_truncated_and_marked():
    long_name = "Dietary change - " + "sustainability reasons " * 20
    fitted = fit_to_tile("[SEGMENT_1_NAME]", long_name)
    assert fitted.endswith(ELLIPSIS)
    assert len(fitted) < len(long_name)


def test_truncated_text_actually_fits_its_tile():
    tile = budget_for("[INSIGHT_1_STAT]")
    fitted = fit_to_tile("[INSIGHT_1_STAT]", "of home bakers " * 20)
    assert len(rendered_lines(fitted, tile.width, tile.typeface, tile.size)) <= tile.lines


def test_truncation_falls_on_a_word_boundary():
    """What survives is whole words from the front — never half a word."""
    original = "of home bakers plan their bakes around a seasonal moment every year"
    fitted = fit_to_tile("[INSIGHT_1_STAT]", original)

    assert fitted.endswith(ELLIPSIS)
    kept = fitted[: -len(ELLIPSIS)].split()
    assert kept, "everything was cut"
    assert kept == original.split()[: len(kept)]


def test_an_unbreakable_word_is_cut_rather_than_left_to_overflow():
    tile = budget_for("[SEGMENT_1_NAME]")
    fitted = fit_to_tile("[SEGMENT_1_NAME]", "A" * 400)
    assert len(rendered_lines(fitted, tile.width, tile.typeface, tile.size)) <= tile.lines


def test_blank_values_stay_blank():
    assert fit_to_tile("[SEGMENT_1_NAME]", "") == ""
    assert fit_to_tile("[SEGMENT_1_NAME]", None) == ""


def test_a_marker_with_no_budget_passes_through_unchanged():
    """Provenance is a flowing block, not a tile — it is measured as a block by
    tests/test_deck_builder.py and must not be truncated line by line."""
    line = "Recommended Products — Source: " + " ".join(["the benchmarking sheet"] * 20)
    assert fit_to_tile("[PROVENANCE_1]", line) == line


def test_every_indexed_tile_shares_one_budget():
    """Tile 3 must not truncate at a different length from tile 1 — three tiles
    that look identical have to break identically."""
    assert budget_for("[SEGMENT_1_NAME]") == budget_for("[SEGMENT_4_NAME]")
    assert budget_for("[INSIGHT_1_STAT]") == budget_for("[INSIGHT_3_STAT]")
    assert budget_for("[PRODUCT_1_FORMAT]") == budget_for("[PRODUCT_3_FORMAT]")


# ---------------------------------------------------------------------------
# The budgets are the template's own numbers
# ---------------------------------------------------------------------------


def _marker_shapes():
    """Every marker in the template: (slide index, marker, shape, run)."""
    found = []
    for index, slide in enumerate(Presentation(TEMPLATE_PATH).slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for marker in MARKER_RE.findall(run.text):
                        found.append((index, marker, shape, run))
    return found


def test_every_template_marker_has_a_budget_or_a_documented_exemption():
    for _, marker, _, _ in _marker_shapes():
        key = marker_family(marker)
        assert key in TILE_BUDGETS or key in EXEMPT_MARKERS, (
            "%s has no tile budget — a new marker must be measured against the "
            "template before it can be filled" % marker
        )


def test_budgets_use_the_typeface_and_size_the_template_uses():
    """The budget is a reading of the template, so a template restyle that
    changes a tile's font has to fail here rather than silently mis-measure."""
    for _, marker, _, run in _marker_shapes():
        tile = budget_for(marker)
        if tile is None:
            continue
        assert run.font.name == tile.typeface, (
            "%s is %s in the template, %s in TILE_BUDGETS"
            % (marker, run.font.name, tile.typeface)
        )
        assert run.font.size.pt == tile.size, (
            "%s is %spt in the template, %spt in TILE_BUDGETS"
            % (marker, run.font.size.pt, tile.size)
        )


def test_every_budget_width_is_exactly_its_box_width():
    """The budget must be the width PowerPoint will actually wrap at.

    Anything narrower is a fiction — the renderer wraps at the box, not at our
    preference, so a value fitted to a narrower budget still draws across the
    neighbouring tile. That was the second half of #162, and it is why
    scripts/narrow_overwide_tiles.py exists: the boxes were brought in to their
    columns rather than the budgets being written down smaller.
    """
    for _, marker, shape, _ in _marker_shapes():
        tile = budget_for(marker)
        if tile is None:
            continue
        usable = wrapping_width(shape)
        # A hundredth of an inch: the budgets are round numbers by choice, and
        # scripts/narrow_overwide_tiles.py leaves a box alone within the same
        # margin. Nothing at this scale is visible.
        assert tile.width == pytest.approx(usable, abs=0.01), (
            "%s is budgeted %.2f\" but its box wraps at %.2f\" — re-run "
            "scripts/narrow_overwide_tiles.py or update TILE_BUDGETS"
            % (marker, tile.width, usable)
        )


def _budget_rects():
    """The rect each tile's text may occupy, at its full budget."""
    rects = {}
    for _, marker, shape, _ in _marker_shapes():
        tile = budget_for(marker)
        if tile is None:
            continue
        left, top = text_origin(shape)
        rects[marker] = (left, top, left + tile.width,
                         top + tile.lines * line_height(tile.typeface, tile.size))
    return rects


def test_marker_text_is_left_aligned_and_top_anchored():
    """The claim that narrowing a box is a no-op for text that already fitted
    rests on this, and on nothing else.

    Left-aligned, top-anchored text starts in the same place whatever the box
    width is, so bringing the right edge in cannot move it. Centre a tile in a
    template revision and that stops being true — every short value shifts —
    so scripts/narrow_overwide_tiles.py would need re-thinking rather than
    re-running.
    """
    for _, marker, shape, _ in _marker_shapes():
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                continue
            assert para.alignment in (None, PP_ALIGN.LEFT), (
                "%s is %s-aligned; the budgets assume left" % (marker, para.alignment)
            )
        anchor = shape.text_frame.vertical_anchor
        assert anchor in (None, MSO_ANCHOR.TOP), (
            "%s is anchored %s; the budgets assume text grows downward from the "
            "top of the box" % (marker, anchor)
        )


def test_no_two_tile_budgets_overlap_each_other():
    """Two tiles filled to the brim must still not touch — this is the whole
    point of #162."""
    rects = _budget_rects()
    slides = {}
    for _, marker, shape, _ in _marker_shapes():
        slides.setdefault(id(shape.part), []).append(marker)

    for markers in slides.values():
        for i, a in enumerate(markers):
            for b in markers[i + 1:]:
                if a not in rects or b not in rects or a == b:
                    continue
                assert not overlaps(rects[a], rects[b]), (
                    "tiles %s and %s overlap when both are full" % (a, b)
                )


def test_no_tile_budget_runs_off_the_slide_or_over_the_artwork():
    """The Prism and Immediate logos sit on the lower half of several slides;
    a full tile must not reach them either."""
    prs = Presentation(TEMPLATE_PATH)
    rects = _budget_rects()
    width, height = inches(prs.slide_width), inches(prs.slide_height)

    for slide in prs.slides:
        markers = [m for shape in slide.shapes if shape.has_text_frame
                   for m in MARKER_RE.findall(shape.text_frame.text)]
        logos = artwork_rects(slide, width)
        for marker in markers:
            rect = rects.get(marker)
            if rect is None:
                continue
            assert rect[2] <= width and rect[3] <= height, (
                "%s runs off the slide when full" % marker
            )
            for _, logo in logos:
                assert not overlaps(rect, logo), (
                    "%s reaches the logo artwork when full" % marker
                )


# ---------------------------------------------------------------------------
# The budgets are big enough for the real data
# ---------------------------------------------------------------------------


def _column(path, field):
    with open(os.path.join(PROJECT_ROOT, path)) as f:
        return [row[field] for row in csv.DictReader(f)]


def test_every_format_name_in_the_catalogue_fits_untruncated():
    """A truncated product name on a client deck is a defect of its own. The
    catalogue's longest name has to fit the tile as authored."""
    truncated = [name for name in set(_column("data/format_recommendations.csv", "Format"))
                 if fit_to_tile("[PRODUCT_1_FORMAT]", name) != name]
    assert truncated == [], truncated


def test_every_reach_figure_fits_untruncated():
    """A truncated number would be a wrong number."""
    reaches = ["{:,}".format(int(r)) for r in _column("data/segments.csv", "reach")]
    truncated = [r for r in set(reaches) if fit_to_tile("[SEGMENT_1_REACH]", r) != r]
    assert truncated == [], truncated


def test_every_segment_name_in_the_catalogue_fits_untruncated():
    """Segment names are the one field the catalogue lets run long (to 95
    characters). The tile takes four lines so the longest still lands whole."""
    names = set(_column("data/segments.csv", "segment_name"))
    truncated = [n for n in names if fit_to_tile("[SEGMENT_1_NAME]", n) != n]
    assert truncated == [], truncated
