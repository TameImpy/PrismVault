"""
Tests for the deck builder at src/deck_builder.py.

No mocks — feeds in a structured brief payload and inspects the .pptx output
against the official 8-slide template (PRD #131 / slice #134). The assertions
are on externally observable properties of the generated file: which markers
were filled, whether the values match the source payload verbatim, whether each
populated run kept the template's own formatting, and whether the static
boilerplate slides came through untouched.
"""
import sys
import os
import io
import math
import re
import zipfile

PROJECT_ROOT = os.path.join(os.path.dirname(__file__), "..")
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from src.deck_builder import TEMPLATE_PATH, build_deck

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_MARKER_RE = re.compile(r"\[[A-Z0-9_]+\]")

# Static boilerplate slides — never populated, must survive byte-identical.
# The Appendix (6) left this list when it gained the provenance markers (#156).
STATIC_SLIDES = (1, 7)

# Structured rows exactly as src.formats.load_format_rows() emits them.
FORMAT_ROWS = [
    {"format": "Infinity Skin", "format_family": "High impact", "ctr": "0.42%",
     "viewability": "78.30%", "indicative_cost": "High", "primary_objective": "Awareness",
     "secondary_objective": "Consideration", "best_for_brief": "Brand launches",
     "best_for_advertiser_type": "Premium brands"},
    {"format": "Playstream Video", "format_family": "Video", "ctr": "0.31%",
     "viewability": "91.04%", "indicative_cost": "Medium", "primary_objective": "Awareness",
     "secondary_objective": "Reach", "best_for_brief": "Video storytelling",
     "best_for_advertiser_type": "Video-led advertisers"},
    {"format": "Standard display - MPU", "format_family": "Standard display", "ctr": "0.05%",
     "viewability": "39.89%", "indicative_cost": "Low", "primary_objective": "Reach",
     "secondary_objective": "Awareness", "best_for_brief": "Efficient reach",
     "best_for_advertiser_type": "Performance-led advertisers"},
]

AUDIENCE_PAYLOAD = {
    "matched": True,
    "note": "Reach figures are per-segment and must not be summed.",
    "query_terms": ["baking"],
    "platforms": [
        {"platform": "AP", "framing": "Modelled first-party audience — broad reach",
         "segments": [
             {"segment_name": "Confident bakers", "reach": 10486296, "platform": "AP",
              "category": "Food & Drink", "plain_english": "Feel confident when baking"},
             {"segment_name": "Cake decorators", "reach": 3401118, "platform": "AP",
              "category": "Food & Drink", "plain_english": "Decorate cakes at home"},
         ]},
        {"platform": "Permutive", "framing": "Behavioural — recently engaged browsers",
         "segments": [
             {"segment_name": "Baking Fans", "reach": 2282660, "platform": "Permutive",
              "category": "Food & Drink", "plain_english": "Regularly browses baking content"},
             {"segment_name": "Dessert Fans", "reach": 941072, "platform": "Permutive",
              "category": "Food & Drink", "plain_english": "Regularly browses dessert content"},
         ]},
    ],
}

INSIGHTS = [
    {"stat": "68%", "text": "of home bakers plan bakes around a season"},
    {"stat": "2.4x", "text": "more likely to try a new brand mid-bake"},
    {"stat": "41%", "text": "read recipe content weekly or more often"},
]

ADVERTISER_OVERVIEW = (
    "Homepride is a heritage UK baking brand best known for its flour and sauces, "
    "currently leaning into everyday family cooking occasions."
)

# Provenance entries exactly as src.provenance.build_provenance() emits them —
# a full brief, so every appendix line has a section behind it.
PROVENANCE = [
    {"section": "Advertiser Overview", "source": "Live public web search",
     "as_at": "2026-08-05", "coverage": "External to the network",
     "period": "As published", "cadence": "Live on every run"},
    {"section": "Client Relationship", "source": "The manual campaign delivery export",
     "as_at": "2026-04-09", "coverage": "Whole network, direct-sold only",
     "period": "2024-01-01 to 2026-01-01", "cadence": "By hand"},
    {"section": "Audience Segments & Reach", "source": "data/segments.csv",
     "as_at": "2026-07-21", "coverage": "Whole network",
     "period": "Rolling 90 days", "cadence": "On request"},
    {"section": "Recommended Products", "source": "The central benchmarking sheet",
     "as_at": "2026-07-20", "coverage": "Whole network",
     "period": "Rolling 12 months", "cadence": "Monthly"},
    {"section": "Historical Research", "source": "Immediate Media — Travel Audience Research",
     "as_at": "2024-12-02", "coverage": "Survey of Immediate audiences (n=1,589)",
     "period": "Fieldwork 2024-11-26 to 2024-12-02", "cadence": "On publication"},
    {"section": "Google Trends", "source": "Google Trends via pytrends",
     "as_at": "2026-08-05", "coverage": "Worldwide search interest",
     "period": "Rolling 12 months", "cadence": "Live on every run"},
]


def _content(**overrides):
    content = {
        "advertiser_overview": ADVERTISER_OVERVIEW,
        "insights": list(INSIGHTS),
        "audience_segments": AUDIENCE_PAYLOAD,
        "format_recommendations": list(FORMAT_ROWS),
        "provenance": list(PROVENANCE),
    }
    content.update(overrides)
    return content


def _slide_text(slide):
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def _deck_text(prs):
    return " ".join(_slide_text(s) for s in prs.slides)


def _find_slide(prs, heading):
    return next(s for s in prs.slides if heading in _slide_text(s))


@pytest.fixture(scope="module")
def deck():
    """A deck built from the full structured payload (research matched)."""
    return Presentation(build_deck(_content(), "Homepride"))


# ---------------------------------------------------------------------------
# Template structure and marker fill
# ---------------------------------------------------------------------------


def test_build_deck_returns_the_eight_slide_template(deck):
    assert len(deck.slides) == 8


def test_build_deck_returns_a_bytesio():
    assert isinstance(build_deck(_content(), "Homepride"), io.BytesIO)


def test_no_marker_survives_in_the_output(deck):
    leftovers = _MARKER_RE.findall(_deck_text(deck))
    assert leftovers == [], "unfilled markers in the generated deck: %s" % leftovers


def test_title_slide_carries_the_advertiser_name(deck):
    assert _slide_text(deck.slides[0]) == "Homepride Brief"


def test_advertiser_overview_prose_lands_on_slide_two(deck):
    text = _slide_text(deck.slides[2])
    assert "Advertiser Overview" in text  # static heading kept
    assert ADVERTISER_OVERVIEW in text


# ---------------------------------------------------------------------------
# Verbatim data transfer — the drift this slice exists to eliminate
# ---------------------------------------------------------------------------


def test_product_metrics_are_verbatim(deck):
    text = _slide_text(_find_slide(deck, "Recommended Products"))
    for row in FORMAT_ROWS:
        assert row["format"] in text
        assert row["ctr"] in text
        assert row["viewability"] in text


def test_segment_names_and_reach_are_verbatim(deck):
    text = _slide_text(_find_slide(deck, "Recommended Segments"))
    for platform in AUDIENCE_PAYLOAD["platforms"]:
        for seg in platform["segments"]:
            assert seg["segment_name"] in text
    # Reach digits are the payload's own — only thousands separators are added.
    assert "10,486,296" in text
    assert "2,282,660" in text


def test_segment_tiles_draw_from_both_platforms():
    """Four tiles, two platforms — the deck must not spend them all on one."""
    prs = Presentation(build_deck(_content(), "Homepride"))
    text = _slide_text(_find_slide(prs, "Recommended Segments"))
    assert "Confident bakers" in text  # AP
    assert "Baking Fans" in text       # Permutive


def test_surplus_tiles_are_blanked_not_left_as_markers():
    """Fewer products than tiles: the spare tiles carry no text at all."""
    prs = Presentation(build_deck(_content(format_recommendations=FORMAT_ROWS[:1]), "Homepride"))
    text = _slide_text(_find_slide(prs, "Recommended Products"))
    assert "Infinity Skin" in text
    assert _MARKER_RE.findall(text) == []
    assert "Playstream Video" not in text


def test_missing_structured_payload_still_builds_a_clean_deck():
    """An older client posting no structured data gets a deck with blank tiles,
    not a crash and not stray markers."""
    prs = Presentation(build_deck({"advertiser_overview": "", "insights": []}, "Homepride"))
    assert _MARKER_RE.findall(_deck_text(prs)) == []


# ---------------------------------------------------------------------------
# Historical research
# ---------------------------------------------------------------------------


def test_historical_insights_are_shown_when_research_matched(deck):
    text = _slide_text(_find_slide(deck, "Historical Insights"))
    for insight in INSIGHTS:
        assert insight["stat"] in text
        assert insight["text"] in text


def test_historical_insights_slide_is_omitted_without_research():
    prs = Presentation(build_deck(_content(insights=[]), "Homepride"))
    assert len(prs.slides) == 7
    assert "Historical Insights" not in _deck_text(prs)
    # Everything else still lands.
    assert "Confident bakers" in _deck_text(prs)


# ---------------------------------------------------------------------------
# Provenance — the deck shares the brief's figures, so it shares their sources
# ---------------------------------------------------------------------------


def test_appendix_names_the_source_of_every_section(deck):
    text = _slide_text(_find_slide(deck, "Appendix"))
    for entry in PROVENANCE:
        assert entry["section"] in text
        assert entry["source"] in text


def test_appendix_carries_all_four_traceability_fields(deck):
    """Source, as-at, coverage and period — the four fields #156 asks for."""
    text = _slide_text(_find_slide(deck, "Appendix"))
    segments = next(e for e in PROVENANCE if e["section"] == "Audience Segments & Reach")
    for value in (segments["source"], segments["as_at"], segments["coverage"],
                  segments["period"]):
        assert value in text


def test_appendix_states_the_benchmark_refresh_cadence(deck):
    """Benchmarks name their origin and how often they are refreshed."""
    text = _slide_text(_find_slide(deck, "Appendix"))
    assert "The central benchmarking sheet" in text
    assert "Monthly" in text


def test_appendix_lines_are_one_per_section_not_one_per_figure(deck):
    """Four product tiles and four segment tiles, one source line each."""
    text = _slide_text(_find_slide(deck, "Appendix"))
    assert text.count("Audience Segments & Reach —") == 1
    assert text.count("Recommended Products —") == 1


def test_appendix_blanks_the_lines_a_brief_did_not_use():
    """A brief with no trends and no matched research leaves no empty labels."""
    prs = Presentation(build_deck(_content(provenance=PROVENANCE[:2]), "Homepride"))
    text = _slide_text(_find_slide(prs, "Appendix"))
    assert "Client Relationship" in text
    assert "Google Trends" not in text
    assert _MARKER_RE.findall(text) == []


def test_the_real_registry_fits_on_the_appendix_slide():
    """The appendix has to hold every source line without running off the slide.

    Measured rather than eyeballed, because the failure mode is a later, wholly
    reasonable edit to `data/provenance.csv` prose — and overflowing deck text
    is a defect this same QA round already raised (#162). The estimate is
    deliberately conservative: average glyph width is taken as half the point
    size, which over-counts for Barlow.
    """
    from src.provenance import build_provenance

    prs = Presentation(build_deck(
        _content(provenance=build_provenance(historical_research={
            "relevant": True, "title": "Immediate Media Travel Audience Research",
            "organisation": "Immediate Media", "fieldwork": "2024-11-26 to 2024-12-02",
            "total_respondents": 1589,
        })), "Homepride"))

    slide = _find_slide(prs, "Appendix")
    box = next(s for s in slide.shapes
               if s.has_text_frame and "[PROVENANCE" not in s.text_frame.text
               and "sources" in s.text_frame.text)

    point_size = box.text_frame.paragraphs[1].runs[0].font.size.pt
    chars_per_line = (box.width / 914400.0) / (point_size / 2.0 / 72.0)
    line_height = point_size * 1.2 / 72.0  # inches

    lines = sum(max(1, math.ceil(len(p.text) / chars_per_line))
                for p in box.text_frame.paragraphs)
    bottom = (box.top / 914400.0) + lines * line_height
    slide_height = prs.slide_height / 914400.0

    assert bottom <= slide_height, (
        "the appendix source lines need %.2f\" and the slide is %.2f\" — trim "
        "data/provenance.csv" % (bottom, slide_height)
    )


def test_appendix_survives_a_deck_built_without_provenance():
    """An older client posts no provenance: the appendix is bare, not broken."""
    prs = Presentation(build_deck(_content(provenance=None), "Homepride"))
    text = _slide_text(_find_slide(prs, "Appendix"))
    assert "Appendix" in text
    assert _MARKER_RE.findall(text) == []


# ---------------------------------------------------------------------------
# Fidelity — styling comes from the template, nothing is hardcoded
# ---------------------------------------------------------------------------


def _run_coords(prs):
    """Map (slide, shape, paragraph, run) -> (text, serialised rPr)."""
    runs = {}
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    rPr = run._r.find(qn("a:rPr"))
                    runs[(si, shi, pi, ri)] = (
                        run.text,
                        etree.tostring(rPr) if rPr is not None else None,
                    )
    return runs


def test_populated_runs_keep_the_template_run_formatting(deck):
    """Every run that held a marker in the template carries byte-identical run
    properties (font, size, colour, spacing) in the output."""
    template_runs = _run_coords(Presentation(TEMPLATE_PATH))
    output_runs = _run_coords(deck)

    checked = 0
    for coords, (text, rPr) in template_runs.items():
        if not _MARKER_RE.search(text):
            continue
        assert coords in output_runs, "template run at %s vanished from the output" % (coords,)
        assert output_runs[coords][1] == rPr, (
            "run properties changed at %s — styling must come from the template" % (coords,)
        )
        checked += 1
    assert checked >= 20, "expected every marker run to be checked, saw %d" % checked


def test_no_hardcoded_styling_in_the_deck_builder():
    """The build path sets text only: no font, size or colour literals."""
    with open(os.path.join(PROJECT_ROOT, "src", "deck_builder.py")) as f:
        source = f.read()
    for banned in ("RGBColor", "Pt(", "font.name", "font.size", "font.bold", "font.color"):
        assert banned not in source, "deck_builder must not hardcode styling (%s)" % banned


def test_marker_split_across_runs_is_still_filled():
    """A template revision can split a marker across runs (PowerPoint does this
    after an edit). The fill still lands, in the first run's formatting."""
    from copy import deepcopy
    from src.deck_builder import _fill_markers

    prs = Presentation(TEMPLATE_PATH)
    para = prs.slides[0].shapes[1].text_frame.paragraphs[0]
    head = para.runs[0]
    tail = deepcopy(head._r)
    head._r.addnext(tail)
    head.text = "[ADVER"
    tail.find(qn("a:t")).text = "TISER_NAME] Brief"

    _fill_markers(prs, {"[ADVERTISER_NAME]": "Homepride"})

    assert _slide_text(prs.slides[0]) == "Homepride Brief"


def test_unknown_marker_is_blanked_rather_than_shipped():
    """A marker the code doesn't know about (added by a template revision) must
    not reach a client — it blanks, and the mapping gap shows as a gap."""
    from src.deck_builder import _fill_markers

    prs = Presentation(TEMPLATE_PATH)
    _fill_markers(prs, {})

    assert _MARKER_RE.findall(_deck_text(prs)) == []


def test_static_boilerplate_slides_are_unchanged(deck):
    """Prism Overview, Appendix and Methodology come through untouched."""
    template = Presentation(TEMPLATE_PATH)
    for index in STATIC_SLIDES:
        expected = etree.tostring(template.slides[index]._element)
        actual = etree.tostring(deck.slides[index]._element)
        assert actual == expected, "static slide %d was modified" % index


def test_brand_fonts_are_embedded():
    buf = build_deck(_content(), "Homepride")
    buf.seek(0)
    with zipfile.ZipFile(buf) as zf:
        names = zf.namelist()
        assert [n for n in names if n.startswith("ppt/fonts/") and n.endswith(".fntdata")]
        assert b"embeddedFontLst" in zf.read("ppt/presentation.xml")


# ---------------------------------------------------------------------------
# File integrity — the deck must open without a repair prompt
# ---------------------------------------------------------------------------


def _slide_roots(buf):
    buf.seek(0)
    with zipfile.ZipFile(buf) as zf:
        for name in sorted(zf.namelist()):
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                yield name, etree.fromstring(zf.read(name))


def _empty_run_count(roots):
    """Runs whose <a:t> is missing or empty — PowerPoint reads those as corrupt."""
    count = 0
    for _, root in roots:
        for r in root.iter("{%s}r" % _A_NS):
            t = r.find("{%s}t" % _A_NS)
            if t is None or t.text is None:
                count += 1
    return count


def test_generated_deck_has_no_ooxml_corruption():
    errors = []
    for name, root in _slide_roots(build_deck(_content(), "Homepride")):
        for t in root.iter("{%s}t" % _A_NS):
            if t.text and "\n" in t.text:
                errors.append("%s: literal newline in <a:t>" % name)
        for body in root.iter("{%s}bodyPr" % _A_NS):
            autofits = [x for x in ("spAutoFit", "normAutofit", "noAutofit")
                        if body.find("{%s}%s" % (_A_NS, x)) is not None]
            if len(autofits) > 1:
                errors.append("%s: conflicting autofit" % name)
    assert errors == [], errors


def test_filling_never_leaves_an_empty_run_behind():
    """The template carries a handful of empty runs of its own (PowerPoint's
    own output). The fill must not add any — a blanked tile drops its run
    rather than leaving <a:t/>."""
    with open(TEMPLATE_PATH, "rb") as f:
        baseline = _empty_run_count(_slide_roots(io.BytesIO(f.read())))

    built = _empty_run_count(_slide_roots(build_deck(_content(), "Homepride")))
    blanked = _empty_run_count(_slide_roots(
        build_deck({"advertiser_overview": "", "insights": []}, "Homepride")))

    assert built <= baseline
    assert blanked <= baseline


def test_multiline_prose_never_reaches_an_at_element():
    """Prose arriving with newlines is flattened rather than embedded raw."""
    content = _content(advertiser_overview="First line.\nSecond line.")
    for name, root in _slide_roots(build_deck(content, "Homepride")):
        for t in root.iter("{%s}t" % _A_NS):
            assert "\n" not in (t.text or "")
    prs = Presentation(build_deck(content, "Homepride"))
    assert "First line. Second line." in _slide_text(prs.slides[2])


def test_dropping_the_insight_slide_leaves_no_dangling_relationship():
    """Omitting the research slide must not leave an orphan slide part."""
    buf = build_deck(_content(insights=[]), "Homepride")
    buf.seek(0)
    with zipfile.ZipFile(buf) as zf:
        slides = [n for n in zf.namelist()
                  if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        assert len(slides) == 7
        rels = etree.fromstring(zf.read("ppt/_rels/presentation.xml.rels"))
        targets = [r.get("Target") for r in rels]
        slide_targets = [t for t in targets if t.startswith("slides/")]
        assert len(slide_targets) == 7


# ---------------------------------------------------------------------------
# API endpoint
# ---------------------------------------------------------------------------

from unittest.mock import patch
from fastapi.testclient import TestClient
from api.main import app

RESEARCH_PAYLOAD = {
    "relevant": True,
    "file": "travel_audience_research_2024_25.md",
    "title": "Travel Audience Research",
    "organisation": "Prism",
    "fieldwork": "2024 to 2025",
    "total_respondents": 2000,
    "matched_on": ["city breaks"],
}


@pytest.fixture()
def api_client(db_url):
    """TestClient backed by the db_url fixture for test isolation."""
    yield TestClient(app)


def test_api_download_deck_returns_401_without_auth(api_client):
    resp = TestClient(app, cookies={}).post(
        "/api/download-deck",
        json={"content": "brief", "topic": "t", "advertiser": "a", "kpi": "k"},
    )
    assert resp.status_code == 401


def _post_deck(api_client, mock_build, **payload):
    """Sign in, post a deck request, and return (response, slide_content)."""
    api_client.post(
        "/api/auth/signup",
        json={"email": "payload@immediate.co.uk", "name": "Test", "password": "password123"},
    )
    mock_build.return_value = io.BytesIO(b"pptx")
    body = {"content": "brief", "topic": "city breaks", "advertiser": "Yakult", "kpi": "Awareness"}
    body.update(payload)
    resp = api_client.post("/api/download-deck", json=body)
    return resp, (mock_build.call_args[0][0] if mock_build.call_args else None)


@patch("api.main.generate_slide_content")
@patch("api.main.build_deck")
def test_api_download_deck_returns_pptx(mock_build, mock_gen, api_client):
    mock_gen.return_value = _content()
    mock_build.return_value = build_deck(_content(), "Yakult")

    api_client.post(
        "/api/auth/signup",
        json={"email": "test@immediate.co.uk", "name": "Test", "password": "password123"},
    )
    resp = api_client.post(
        "/api/download-deck",
        json={"content": "brief", "topic": "gut health", "advertiser": "Yakult", "kpi": "Awareness"},
    )

    assert resp.status_code == 200
    assert "openxmlformats" in resp.headers["content-type"]
    assert "Prism_Plan_Yakult_gut_health.pptx" in resp.headers["content-disposition"]


@patch("api.main.generate_slide_content")
@patch("api.main.build_deck")
def test_api_download_deck_forwards_structured_payload(mock_build, mock_gen, api_client):
    """All three structured fields reach build_deck exactly as the brief run
    produced them — reach, CTR and viewability never pass back through the LLM,
    so they cannot drift from the brief."""
    from src.formats import load_format_rows

    mock_gen.return_value = _content()
    format_rows = load_format_rows()[:2]

    resp, slide_content = _post_deck(
        api_client, mock_build,
        audience_segments=AUDIENCE_PAYLOAD,
        format_recommendations=format_rows,
        historical_research=RESEARCH_PAYLOAD,
    )

    assert resp.status_code == 200
    assert slide_content["audience_segments"] == AUDIENCE_PAYLOAD
    assert slide_content["historical_research"] == RESEARCH_PAYLOAD
    assert slide_content["format_recommendations"] == format_rows
    assert slide_content["format_recommendations"][0]["ctr"] == format_rows[0]["ctr"]
    assert slide_content["audience_segments"]["platforms"][0]["segments"][0]["reach"] == 10486296


@patch("api.main.generate_slide_content")
@patch("api.main.build_deck")
def test_api_download_deck_structured_fields_are_optional(mock_build, mock_gen, api_client):
    """Older clients that post none of the new fields still get a deck."""
    mock_gen.return_value = _content()

    resp, slide_content = _post_deck(api_client, mock_build)

    assert resp.status_code == 200
    assert slide_content["format_recommendations"] is None
    assert slide_content["historical_research"] is None


@patch("api.main.generate_slide_content")
@patch("api.main.build_deck")
def test_api_passes_matched_research_to_the_content_step(mock_build, mock_gen, api_client):
    """Insight selection is grounded in the matched research the brief found."""
    mock_gen.return_value = _content()

    _post_deck(api_client, mock_build, historical_research=RESEARCH_PAYLOAD)

    assert mock_gen.call_args[1]["historical_research"] == RESEARCH_PAYLOAD
